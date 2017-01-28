# -*- coding: utf-8 -*-

"""
lottery_presentation/generate_presentation.py

Generates a ppt presentation from csv dump of admission lottery results.
Call with two positional args; eg.:
    $ python generate_presentation.py <lottery results csv file> <desired output file name>

Assumes the csv is ordered by status, with accepted students first, followed by
waitlist students.  Column headers expected in the csv are:
- id
- lottery_number: will be 'Offered' if student has been accepted, 'WL $xxxx' if student is on wait list
- first_name
- last_name
- Elementary

"""

import argparse
import csv
from queue import Queue

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches


class PresentationMaker():
    """
    A a manager object used to read data from a csv
    of lottery results and add them to a pptx.Presentation object.

    Instantiated with two parameters:
        infile_name:  name of csv lottery results file. Expects the following
                      column headers:
                            * id
                            * lottery_number
                            * first_name
                            * last_name
                            * Elementary

        outfile_name: filename to save the ppt presentation as.
    """

    ## presentation slide spec
    TEMPLATE_FILENAME = "templates/template.pptx"
    TITLE_SLIDE_LAYOUT_INDEX = 0 # title text in middle of slide
    BODY_SLIDE_LAYOUT_INDEX = 5  # empty body section

    ## body slide table spec
    MAX_ROWS     = 8 # in addition to header row on a body slide
    TOP_MARGIN   = Inches(1.7)
    LEFT_MARGIN  = Inches(0)
    TABLE_WIDTH  = Inches(10.0)
    TABLE_HEIGHT = Inches(0.8) # base height; auto adjusted
    COL_1_WIDTH  = Inches(3.0)
    COL_2_WIDTH  = Inches(2.0)
    COL_3_WIDTH  = Inches(5.0)
    NAME_FORMAT  = "{last_name}, {first_name}"
    NOBLE_NAVY   = RGBColor(0, 67, 120)


    def __init__(self, infile_name, outfile_name):
        self.infile_name  = infile_name
        self.outfile_name = outfile_name
        self.presentation = Presentation(self.TEMPLATE_FILENAME)
        self.title_layout = self.presentation.slide_layouts[self.TITLE_SLIDE_LAYOUT_INDEX]
        self.body_layout  = self.presentation.slide_layouts[self.BODY_SLIDE_LAYOUT_INDEX]
        self.body_queue   = Queue()


    def make_presentation(self):
        """
        Read in data from the infile csv, and save a pptx
        under the outfile name.

        Creates two sections (a "section" being a title slide and set of
        following body slides) -- one for admitted students,
        and one for waitlist students.
        """

        self._add_title_slide("Admitted Students")

        with open(self.infile_name) as csvfile:

            reader = csv.DictReader(csvfile)

            # process enrolled slides
            for row in reader:
                if row['lottery_number'] == "Offered":
                    self._add_to_body_queue(row)
                else:
                    # at some point switches to waitlist students
                    self._end_body_section()
                    self._add_title_slide("Waitlist Students")
                    self._add_to_body_queue(row)
                    break
            # process the waitlist students
            for row in reader:
                self._add_to_body_queue(row)

        self._end_body_section()
        self.presentation.save(self.outfile_name)


    def _add_title_slide(self, title_string=''):
        """
        Helper function that adds a title slide with title_string text.
        """

        slide = self.presentation.slides.add_slide(self.title_layout)
        slide.shapes.title.text = title_string


    def _add_to_body_queue(self, row_dict):
        """
        Helper function that adds a row to the body_queue,
        and calls _add_body_slide when body_queue size reaches MAX_ROWS.

        :param row_dict: A row dictionary, as created by :class:`csv.DictReader`
        """

        self.body_queue.put(row_dict)
        if self.body_queue.qsize() >= self.MAX_ROWS:
            self._add_body_slide()


    def _add_body_slide(self):
        """
        Helper function that adds a body slide to the presentation, using
        the rows in the body_queue to fill out the table.
        """

        body_slide = self.presentation.slides.add_slide(self.body_layout)
        shapes = body_slide.shapes

        num_cols = 3
        num_rows = self.MAX_ROWS + 1 # allow for header row

        table = shapes.add_table(
            num_rows,
            num_cols,
            self.LEFT_MARGIN,
            self.TOP_MARGIN,
            self.TABLE_WIDTH,
            self.TABLE_HEIGHT
        ).table

        table.columns[0].width = self.COL_1_WIDTH
        table.columns[1].width = self.COL_2_WIDTH
        table.columns[2].width = self.COL_3_WIDTH

        self._format_table_headers(table)

        row_number = 1 # start after header row
        while self.body_queue.qsize():
            row_dict = self.body_queue.get()
            # table.cell(row,col)
            table.cell(row_number, 0).text = self.NAME_FORMAT.format(**row_dict)
            table.cell(row_number, 1).text = "{lottery_number}".format(**row_dict)
            table.cell(row_number, 2).text = "{Elementary}".format(**row_dict)
            row_number += 1


    def _format_table_headers(self, table):
        """
        Helper function to format table headers.
        """

        for i in range(3):
            cell = table.cell(0, i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.NOBLE_NAVY

        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Status"
        table.cell(0, 2).text = "School"


    def _end_body_section(self):
        """
        Helper function that calls _add_body_slide if at least one row in the
        body_queue, used to end the different body sections of the presentation
        (enrolled and waitlist student sections).
        """

        if self.body_queue.qsize():
            self._add_body_slide()


def parse_args():
    """
    Setup the input and output arguments for the script.
    Return the parsed input and output files.
    """

    parser = argparse.ArgumentParser(description=\
        "Create a PowerPoint presentation from a csv of lottery results"
    )
    parser.add_argument('infile',
                        type=argparse.FileType('r'),
                        help='CSV file to read in'
    )
    parser.add_argument('outfile',
                        type=argparse.FileType('w'),
                        help='Output powerpoint'
    )
    return parser.parse_args()


if __name__ == "__main__":

    args = parse_args()
    presentation = PresentationMaker(args.infile.name, args.outfile.name)
    presentation.make_presentation()
